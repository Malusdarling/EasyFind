"""
reader.py — 文件读取模块
======================
核心功能：多格式文件内容读取、编码自动检测、二进制文件判断

支持的格式：
  - 纯文本类: .txt, .md, .csv, .json, .xml, .ini, .log, .yaml, .toml
  - 代码类: .py, .java, .c, .cpp, .h, .js, .ts, .html, .css, .php, .rb, .go, .rs, .swift
  - Office类: .docx, .xlsx, .pptx（需安装对应库）
  - 其他: 通过编码探测自动尝试读取

编码检测策略（自动回退链）:
  utf-8 → gbk → gb2312 → gb18030 → big5 → latin-1（兜底）
"""

import os
import sys
import struct
from pathlib import Path

# ---------- 编码探测优先级 ----------
# 对于中文环境，优先尝试 UTF-8 和 GB 系列编码
ENCODING_PRIORITY = ['utf-8', 'gbk', 'gb2312', 'gb18030', 'big5', 'latin-1']

# ---------- 用户文档扩展名白名单 ----------
# 仅包含用户可能主动编辑的文档类型，排除代码文件（.py/.java/.js 等）
# 这样在开发目录中不会误扫大量代码文件，大幅提升检索效率
USER_FILE_EXTENSIONS = {
    # 可内容搜索的用户文档
    '.txt',        # 文本
    '.doc', '.docx',   # Word（旧/新格式均支持）
    '.xls', '.xlsx',   # Excel（旧/新格式均支持）
    '.ppt', '.pptx',   # PPT（旧/新格式均支持）
    '.csv',        # 数据
}

# ---------- Office 文档扩展名（新格式专用解析器）----------
OFFICE_EXTENSIONS = {'.docx', '.xlsx', '.pptx'}

# ---------- 旧版 Office 扩展名（二进制格式，olefile/xlrd 解析）----------
LEGACY_OFFICE_EXTENSIONS = {'.doc', '.xls', '.ppt'}

# ---------- 最大可读取文件大小（10MB，防止大文件拖慢）----------
MAX_FILE_SIZE = 10 * 1024 * 1024


# ===================================================================
# 离线/云盘文件检测
# ===================================================================
# 云盘离线文件检测（基于 Windows 属性标志）
# ===================================================================
# WPS 云盘、OneDrive 等同步盘会将未同步到本地的文件标记为
# FILE_ATTRIBUTE_OFFLINE (0x1000)。已同步到本地的文件不会有此标志。
#
# 策略：仅通过 Windows 属性检测，不靠路径关键词判断。
# 已同步到本地的 OneDrive/WPS 文件可以正常搜索。
# 仅过滤真正离线的云占位文件。


def is_offline_cloud_file(filepath: str) -> bool:
    """
    检测文件是否为云盘离线占位文件（不在本地、需要从云端拉取）。

    使用 Windows FILE_ATTRIBUTE_OFFLINE 标志判断：
    - OneDrive / WPS 云盘 / Dropbox 等同步盘对未同步的文件设此标志
    - 已同步到本地的文件不会有此标志，可以正常搜索

    参数:
        filepath: 文件完整路径

    返回:
        True  → 确定是云盘离线占位文件（应跳过搜索）
        False → 本地文件 / 无法判断（应参与搜索）
    """
    if sys.platform != 'win32':
        return False

    try:
        import ctypes
        FILE_ATTRIBUTE_OFFLINE = 0x1000
        attrs = ctypes.windll.kernel32.GetFileAttributesW(str(filepath))
        # 仅当 Windows 明确标记为 OFFLINE 时返回 True
        # NORMAL 文件或其他属性不匹配，返回 False
        return attrs != -1 and bool(attrs & FILE_ATTRIBUTE_OFFLINE)
    except Exception:
        return False  # API 调用失败也视为本地（不让用户丢结果）


def is_file_accessible(filepath: str) -> tuple[bool, str]:
    """
    检查文件是否可访问，返回 (可访问, 原因)。

    原因说明:
        'ok'          — 可正常访问
        'offline'     — 云盘离线文件（不在本地）
        'not_found'   — 文件不存在
        'permission'  — 无权限访问
        'too_large'   — 文件过大
        'unknown'     — 其他原因
    """
    if not os.path.exists(filepath):
        return False, 'not_found'

    if is_offline_cloud_file(filepath):
        return False, 'offline'

    try:
        with open(filepath, 'rb') as f:
            f.read(1)
        return True, 'ok'
    except FileNotFoundError:
        # 路径存在但打开失败——不直接判离线，先查一下 Windows 属性
        if is_offline_cloud_file(filepath):
            return False, 'offline'
        return False, 'not_found'
    except PermissionError:
        return False, 'permission'
    except OSError:
        if is_offline_cloud_file(filepath):
            return False, 'offline'
        return False, 'unknown'
    except Exception:
        return False, 'unknown'


# ===================================================================
# 核心函数：读取文件内容
# ===================================================================
def read_file_content(filepath: str) -> tuple[str | None, str | None]:
    """
    读取指定文件的内容，自动检测编码。

    参数:
        filepath: 文件完整路径

    返回:
        (content, encoding) 元组:
            content  — 文件文本内容（读取失败为 None）
            encoding — 实际使用的编码（读取失败为 None）

    工作流程:
        1. 根据扩展名选择读取策略
        2. 文本文件 → 编码回退链尝试
        3. Office 文件 → 专用解析器
        4. 未知类型 → 二进制检测后尝试文本读取
    """
    ext = Path(filepath).suffix.lower()

    # ----- 新式 Office 文档（OpenXML 格式）-----
    if ext == '.docx':
        return _read_docx(filepath)
    elif ext == '.xlsx':
        return _read_xlsx(filepath)
    elif ext == '.pptx':
        return _read_pptx(filepath)

    # ----- 旧版 Office 文档（二进制格式）-----
    if ext == '.doc':
        return _read_doc(filepath)
    elif ext == '.xls':
        return _read_xls(filepath)
    elif ext == '.ppt':
        return _read_ppt(filepath)

    # ----- 文本文件走编码探测 -----
    if ext in USER_FILE_EXTENSIONS:
        accessible, reason = is_file_accessible(filepath)
        if not accessible:
            return _offline_error_msg(reason, 'txt'), None
        return _read_text_with_encoding_detection(filepath)

    # 不在用户文档类型中 → 不可读
    return None, None


def _read_text_with_encoding_detection(filepath: str) -> tuple[str | None, str | None]:
    """
    使用编码回退链读取文本文件。
    按优先级逐个尝试编码，第一个成功解码的编码即为实际编码。
    """
    for enc in ENCODING_PRIORITY:
        try:
            with open(filepath, 'r', encoding=enc, errors='strict') as f:
                content = f.read()
            return content, enc
        except (UnicodeDecodeError, UnicodeError):
            continue
        except Exception:
            # 文件不存在、权限错误等 → 不再重试
            return None, None

    # 所有编码都失败，用 latin-1（不抛异常，但可能有乱码）
    try:
        with open(filepath, 'r', encoding='latin-1', errors='replace') as f:
            content = f.read()
        return content, 'latin-1'
    except Exception:
        return None, None


# ===================================================================
# 二进制文件检测
# ===================================================================
def _is_binary_file(filepath: str, check_bytes: int = 8192) -> bool:
    """
    通过检查文件头字节判断是否为二进制文件。
    读取文件开头的一段字节，如果包含 NULL 字节则判定为二进制。

    参数:
        filepath: 文件路径
        check_bytes: 检查的字节数（默认 8KB）

    返回:
        True  → 判定为二进制文件
        False → 判定为文本文件
    """
    try:
        with open(filepath, 'rb') as f:
            chunk = f.read(check_bytes)
        # NULL 字节是二进制文件的强信号
        if b'\x00' in chunk:
            return True
        # 检查是否有大量非 ASCII、非常见文本控制字符
        text_characters = (
            b'\n\r\t\b\f' +
            bytes(range(32, 127))  # 可打印 ASCII
        )
        non_text = sum(1 for byte in chunk if byte not in text_characters and byte < 128)
        # 如果非文本字节超过 30%，判定为二进制
        if len(chunk) > 0 and non_text / len(chunk) > 0.3:
            return True
        return False
    except Exception:
        return True  # 无法读取 → 保守视为二进制


# ===================================================================
# 文件大小检查
# ===================================================================
def is_file_too_large(filepath: str) -> bool:
    """
    检查文件是否超过大小限制。
    超大文件跳过内容搜索以免内存溢出。
    """
    try:
        return os.path.getsize(filepath) > MAX_FILE_SIZE
    except Exception:
        return True


# ===================================================================
# Office 文档解析器（可选依赖，缺失时优雅降级）
# ===================================================================
def _read_docx(filepath: str) -> tuple[str | None, str | None]:
    """读取 Word .docx 文件（让 python-docx 自行处理 I/O 错误）"""
    try:
        from docx import Document
        doc = Document(filepath)
        lines = []
        for para in doc.paragraphs:
            lines.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                lines.append(' | '.join(row_text))
        content = '\n'.join(lines)
        return content, 'docx'
    except ImportError:
        return None, None
    except Exception as e:
        error_msg = _parse_office_error(e, filepath, 'docx')
        return error_msg, None


def _read_xlsx(filepath: str) -> tuple[str | None, str | None]:
    """读取 Excel .xlsx 文件（让 openpyxl 自行处理 I/O 错误）"""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(filepath, read_only=True, data_only=True)
        lines = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            lines.append(f"=== 工作表: {sheet_name} ===")
            for row in ws.iter_rows(values_only=True):
                row_str = ' | '.join(str(cell) if cell is not None else '' for cell in row)
                if row_str.strip():
                    lines.append(row_str)
        content = '\n'.join(lines)
        return content, 'xlsx'
    except ImportError:
        return None, None
    except Exception as e:
        error_msg = _parse_office_error(e, filepath, 'xlsx')
        return error_msg, None


def _read_pptx(filepath: str) -> tuple[str | None, str | None]:
    """读取 PowerPoint .pptx 文件（让 python-pptx 自行处理 I/O 错误）"""
    try:
        from pptx import Presentation
        prs = Presentation(filepath)
        lines = []
        for slide_num, slide in enumerate(prs.slides, 1):
            lines.append(f"=== 幻灯片 {slide_num} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        lines.append(para.text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        row_text = [cell.text for cell in row.cells]
                        lines.append(' | '.join(row_text))
        content = '\n'.join(lines)
        return content, 'pptx'
    except ImportError:
        return None, None
    except Exception as e:
        error_msg = _parse_office_error(e, filepath, 'pptx')
        return error_msg, None


# ===================================================================
# 旧版 Office 文档解析器（olefile + xlrd，二进制格式）
# ===================================================================
def _read_doc(filepath: str) -> tuple[str | None, str | None]:
    """
    读取 Word .doc 文件（97-2003 二进制格式）。

    策略（三级降级）：
      1. win32com → 通过 Word COM 接口提取（最准确，需安装 Office/WPS）
      2. olefile → 从 OLE 流中提取可读文本片段
      3. 原始字节 → 直接扫描二进制中的可读字符串
    """
    # ----- 第一级：win32com（需安装 Microsoft Office 或 WPS）-----
    try:
        import win32com.client
        word = win32com.client.Dispatch("Word.Application")
        word.Visible = False
        word.DisplayAlerts = False
        import pythoncom
        pythoncom.CoInitialize()
        try:
            doc = word.Documents.Open(filepath, ReadOnly=True, AddToRecentFiles=False)
            text = doc.Content.Text
            doc.Close()
            if text and text.strip():
                return text.strip(), 'doc'
        except Exception:
            pass
        finally:
            try:
                word.Quit()
            except Exception:
                pass
    except ImportError:
        pass
    except Exception:
        pass

    # ----- 第二级：olefile 从 OLE 流中提取文本 -----
    raw = b''
    try:
        import olefile
        ole = olefile.OleFileIO(filepath)
        all_data = []
        for stream_name in ole.listdir():
            try:
                data = ole.openstream(stream_name).read()
                all_data.append(data)
            except Exception:
                pass
        ole.close()
        raw = b''.join(all_data)
    except ImportError:
        pass
    except Exception:
        try:
            with open(filepath, 'rb') as f:
                raw = f.read()
        except Exception:
            pass

    if raw:
        # 尝试 GBK/UTF-8 解码，过滤乱码行
        for enc in ('gbk', 'utf-8', 'latin-1'):
            try:
                text = raw.decode(enc, errors='replace')
                lines = []
                for line in text.split('\n'):
                    reps = line.count('�') + line.count('?')
                    if len(line) > 3 and reps / max(len(line), 1) < 0.3:
                        clean = ''.join(c for c in line if c.isprintable() or c in '\n\r\t')
                        clean = clean.strip()
                        if len(clean) > 3:
                            lines.append(clean)
                if lines:
                    return '\n'.join(lines), enc
            except Exception:
                continue

        # 兜底：ASCII 可读字符串
        result, cur = [], []
        for byte in raw:
            if 32 <= byte < 127 or byte in (10, 13, 9):
                cur.append(chr(byte))
            else:
                if len(cur) > 3:
                    result.append(''.join(cur))
                cur = []
        if len(cur) > 3:
            result.append(''.join(cur))
        if result:
            return '\n'.join(result), 'ascii'

    return None, None


def _read_xls(filepath: str) -> tuple[str | None, str | None]:
    """
    读取 Excel .xls 文件（97-2003 二进制格式）。
    使用 xlrd 解析。
    """
    try:
        import xlrd
        wb = xlrd.open_workbook(filepath)
        lines = []
        for sheet_idx in range(wb.nsheets):
            ws = wb.sheet_by_index(sheet_idx)
            lines.append(f"=== 工作表: {ws.name} ===")
            for row_idx in range(ws.nrows):
                row_vals = []
                for col_idx in range(ws.ncols):
                    cell = ws.cell(row_idx, col_idx)
                    if cell.ctype == xlrd.XL_CELL_EMPTY:
                        row_vals.append('')
                    elif cell.ctype == xlrd.XL_CELL_DATE:
                        row_vals.append(str(cell.value))
                    else:
                        row_vals.append(str(cell.value))
                row_str = ' | '.join(row_vals)
                if row_str.strip():
                    lines.append(row_str)
        wb.release_resources()
        content = '\n'.join(lines)
        return content, 'xls'
    except ImportError:
        return None, None
    except Exception as e:
        error_msg = _parse_office_error(e, filepath, 'xls')
        return error_msg, None


def _read_ppt(filepath: str) -> tuple[str | None, str | None]:
    """
    读取 PowerPoint .ppt 文件（97-2003 二进制格式）。
    与 _read_doc 相同策略：收集 OLE 流数据，用 GBK/UTF-8 分段解码。
    """
    raw = b''
    try:
        import olefile
        ole = olefile.OleFileIO(filepath)
        all_data = []
        for stream_name in ole.listdir():
            try:
                data = ole.openstream(stream_name).read()
                all_data.append(data)
            except Exception:
                pass
        ole.close()
        raw = b''.join(all_data)
    except ImportError:
        return None, None
    except Exception:
        try:
            with open(filepath, 'rb') as f:
                raw = f.read()
        except Exception:
            return None, None

    if not raw:
        return None, None

    for enc in ('gbk', 'utf-8', 'latin-1'):
        try:
            text = raw.decode(enc, errors='replace')
            lines = []
            for line in text.split('\n'):
                replacements = line.count('�') + line.count('?')
                if len(line) > 2 and replacements / max(len(line), 1) < 0.3:
                    clean = ''.join(c for c in line if c.isprintable() or c in '\n\r\t')
                    clean = clean.strip()
                    if len(clean) > 2:
                        lines.append(clean)
            if lines:
                return '\n'.join(lines), enc
        except Exception:
            continue

    # 兜底 ASCII
    result, current = [], []
    for byte in raw:
        if 32 <= byte < 127 or byte in (10, 13, 9):
            current.append(chr(byte))
        else:
            if len(current) > 3:
                result.append(''.join(current))
            current = []
    if len(current) > 3:
        result.append(''.join(current))
    if result:
        return '\n'.join(result), 'ascii'
    return None, None


def _offline_error_msg(reason: str, ext: str) -> str:
    """生成云盘离线文件友好提示"""
    if reason == 'offline':
        return (f'☁️ 文件在云端，未同步到本地\n\n'
                f'该 {ext} 文件位于云盘（可能是 WPS 云盘 / OneDrive），\n'
                f'当前仅存占位符，实际内容不在本机。\n\n'
                f'解决方法：\n'
                f'  1. 在线打开文件，让客户端自动同步到本地\n'
                f'  2. 在云盘目录中右键文件 → 「始终保留在此设备」\n'
                f'  3. 同步完成后重新搜索即可预览内容\n\n'
                f'文件路径: {ext}')
    elif reason == 'permission':
        return f'⛔ 无权限访问该 {ext} 文件'
    elif reason == 'not_found':
        return f'❌ 文件未找到（可能已被移动或删除）'
    else:
        return f'⚠️ 无法读取 {ext} 文件内容'


def _parse_office_error(e: Exception, filepath: str, ext: str) -> str:
    """解析 Office 文件读取异常，返回中文友好提示"""
    error_str = str(e).lower()

    # 检测各类常见错误
    if 'file not found' in error_str or 'no such file' in error_str:
        return f'❌ 文件未找到（可能已被移动或删除）\n\n文件路径: {filepath}'
    if 'permission' in error_str or 'access denied' in error_str:
        return f'⛔ 无权限访问该文件\n\n文件路径: {filepath}'
    if 'bad magic number' in error_str or 'not a' in error_str and 'file' in error_str:
        return f'⚠️ 文件格式损坏或不是有效的 {ext} 文件\n\n文件路径: {filepath}'
    if 'truncated' in error_str or 'unexpected end' in error_str:
        return f'⚠️ 文件不完整（可能未完全同步到本地）\n\n文件路径: {filepath}'

    # 通用错误
    return (f'⚠️ 无法解析 {ext} 文件内容\n'
            f'  原因: {e}\n'
            f'  提示: 尝试用 WPS Office 手动打开该文件，\n'
            f'        确认文件完整后重新搜索。\n\n'
            f'文件路径: {filepath}')


# ===================================================================
# 工具函数
# ===================================================================
def is_searchable_file(filepath: str) -> bool:
    """
    判断文件是否可参与内容搜索。
    v2.0: 仅搜索用户文档类型，跳过代码/系统文件。
    条件：扩展名在 USER_FILE_EXTENSIONS 或 OFFICE_EXTENSIONS 中，且未超大小限制。
    """
    ext = Path(filepath).suffix.lower()
    if ext in OFFICE_EXTENSIONS:
        return True  # Office 文件交给专用解析器
    if ext in USER_FILE_EXTENSIONS:
        return not is_file_too_large(filepath)
    return False  # 非用户文档类型，不参与内容搜索


def get_file_icon(filepath: str) -> str:
    """根据扩展名返回图标标记（用于结果列表）"""
    ext = Path(filepath).suffix.lower()
    if ext in {'.py', '.java', '.c', '.cpp', '.js', '.ts', '.go', '.rs'}:
        return '📄'
    elif ext in {'.docx', '.doc'}:
        return '📝'
    elif ext in {'.xlsx', '.xls'}:
        return '📊'
    elif ext in {'.pptx', '.ppt'}:
        return '📑'
    elif ext in {'.txt', '.md'}:
        return '📃'
    elif ext in {'.pdf'}:
        return '📕'
    elif ext in {'.html', '.htm', '.css', '.js'}:
        return '🌐'
    else:
        return '📄'
